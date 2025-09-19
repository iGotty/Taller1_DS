#!/usr/bin/env python3
"""
Generador de presentación ejecutiva para el análisis de cancelaciones hoteleras
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np
from pathlib import Path

def create_presentation():
    """Genera presentación PowerPoint con resultados del análisis"""
    
    # Crear presentación
    prs = Presentation()
    
    # Configurar dimensiones
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Título
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Análisis de Cancelaciones Hoteleras"
    subtitle.text = "Estrategia de Reducción y Optimización de Ocupación\nSeptiembre 2025"
    
    # SLIDE 2: Contexto
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contexto del Análisis"
    content.text = ("• Dataset: 58,895 reservas hoteleras\n"
                   "• Período: 2015-2017\n"
                   "• Tipos: City Hotel y Resort Hotel\n"
                   "• Objetivo: Reducir cancelaciones y mejorar ocupación\n"
                   "• Enfoque: Análisis predictivo y segmentación")
    
    # SLIDE 3: Problema
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Problemática Actual"
    content.text = ("• Tasa de cancelación: 41.1%\n"
                   "• 24,224 cancelaciones totales\n"
                   "• Pérdida estimada: $3.6M anuales\n"
                   "• Variabilidad extrema entre segmentos\n"
                   "• Sin políticas diferenciadas actualmente")
    
    # SLIDE 4: Metodología
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Metodología Aplicada"
    content.text = ("1. Análisis Exploratorio\n"
                   "   • 32 variables analizadas\n"
                   "   • Identificación de patrones\n\n"
                   "2. Tests Estadísticos\n"
                   "   • Chi-cuadrado para categóricas\n"
                   "   • Mann-Whitney para numéricas\n\n"
                   "3. Modelado Predictivo\n"
                   "   • Regresión logística\n"
                   "   • ROC-AUC: 0.78")
    
    # SLIDE 5: Hallazgo 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Hallazgo #1: Lead Time como Factor Crítico"
    content.text = ("• Reservas >60 días: 52% cancelación\n"
                   "• Reservas <30 días: 28% cancelación\n"
                   "• Diferencia: 24 puntos porcentuales\n\n"
                   "Implicación:\n"
                   "• Políticas diferenciadas por ventana de reserva\n"
                   "• Mayor riesgo requiere mayor garantía")
    
    # SLIDE 6: Hallazgo 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Hallazgo #2: Impacto de Depósitos"
    content.text = ("Sin depósito: 46.2% cancelación\n"
                   "Con depósito no reembolsable: 4.7% cancelación\n\n"
                   "Reducción: 89% en tasa de cancelación\n\n"
                   "Oportunidad:\n"
                   "• Implementar depósitos obligatorios\n"
                   "• Segmentar por lead time y canal")
    
    # SLIDE 7: Hallazgo 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Hallazgo #3: Variabilidad por Canal"
    content.text = ("TA/TO: 48% cancelación\n"
                   "Direct: 24% cancelación\n"
                   "Corporate: 19% cancelación\n\n"
                   "Estrategia:\n"
                   "• Políticas diferenciadas por canal\n"
                   "• Incentivos para canales directos\n"
                   "• Renegociación con OTAs")
    
    # SLIDE 8: Recomendaciones
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Recomendaciones Priorizadas"
    content.text = ("INMEDIATAS (Semanas 1-4):\n"
                   "1. Depósitos obligatorios para lead_time >60 días\n"
                   "2. Sistema de alertas para alto riesgo\n\n"
                   "CORTO PLAZO (Meses 1-3):\n"
                   "3. Modelo predictivo en producción\n"
                   "4. Contacto proactivo con clientes\n\n"
                   "MEDIANO PLAZO (Meses 3-6):\n"
                   "5. Pricing dinámico\n"
                   "6. Overbooking inteligente")
    
    # SLIDE 9: Plan de Implementación
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Plan de Implementación"
    content.text = ("Fase 1: Quick Wins (Semanas 1-4)\n"
                   "• Políticas de depósito\n"
                   "• Costo: $50K\n\n"
                   "Fase 2: Optimización (Semanas 5-16)\n"
                   "• Modelo predictivo\n"
                   "• Costo: $150K\n\n"
                   "Fase 3: Automatización (Semanas 17-40)\n"
                   "• Sistema completo\n"
                   "• Costo: $300K")
    
    # SLIDE 10: Impacto Proyectado
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Impacto Económico Proyectado"
    content.text = ("Situación actual: 41.1% cancelación\n"
                   "Objetivo: 32.0% cancelación\n"
                   "Reducción: 9.1 puntos porcentuales\n\n"
                   "Beneficios anuales:\n"
                   "• Ingresos recuperados: $1.4M\n"
                   "• Inversión total: $500K\n"
                   "• ROI: 2.8x primer año\n"
                   "• Payback: 4 meses")
    
    # SLIDE 11: Próximos Pasos
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Próximos Pasos"
    content.text = ("Semana 1:\n"
                   "✓ Aprobación ejecutiva\n"
                   "✓ Formación de equipo\n\n"
                   "Semana 2:\n"
                   "✓ Diseño de políticas\n"
                   "✓ Configuración de métricas\n\n"
                   "Mes 1:\n"
                   "✓ Implementación Fase 1\n"
                   "✓ Inicio desarrollo modelo")
    
    # SLIDE 12: Conclusiones
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusiones"
    content.text = ("• Oportunidad clara de mejora identificada\n"
                   "• Factores críticos cuantificados\n"
                   "• Plan de acción concreto y medible\n"
                   "• ROI atractivo con riesgo controlado\n\n"
                   "RECOMENDACIÓN:\n"
                   "Proceder con implementación inmediata\n"
                   "de la estrategia en 3 fases")
    
    # Guardar presentación
    output_path = Path("reports/presentacion_taller.pptx")
    output_path.parent.mkdir(exist_ok=True)
    prs.save(str(output_path))
    
    print(f"Presentación generada exitosamente: {output_path}")
    return str(output_path)

if __name__ == "__main__":
    create_presentation()